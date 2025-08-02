import os
import nibabel as nib
import numpy as np
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime
from tqdm import tqdm
import sys
# ========== CONFIG ==========

folder_path = sys.argv[1]
# folder_path = "/Volumes/pr_ohlendorf/fMRI/AnalysedData/Project_MMP9_NJ_MP/Project1_Active_vs_Inactive/RGRO_250402_0224_RN_SD_018/17functionalEPI"
output_pptx = os.path.join(folder_path, "preprocessing_report.pptx")
temp_dir = os.path.join(folder_path, "temp_imgs")
os.makedirs(temp_dir, exist_ok=True)

file_list = [
    "G1_cp.nii.gz",
    "mc_func.nii.gz",
    "rest_rotation.jpg",
    "rest_translation.jpg",
    "mean_mc_func.nii.gz",
    "initial_mean_mc_func.nii.gz",
    "mask_mean_mc_func.nii.gz",
    "cleaned_mc_func.nii.gz",
    "cleaned_N4_mean_mc_func.nii.gz",
    "before_despiking_spikecountTC.png",
    "after_despiking_spikecountTC.png",
    "despike_cleaned_mc_func.nii.gz",
    "cleaned_sm_despike_cleaned_mc_func.nii.gz"
]

title_map = {
    "G1_cp.nii.gz": "Initial Raw 4D data",
    "mc_func.nii.gz": "Motion Corrected 4D data",
    "rest_rotation.jpg": "Rotational Movement",
    "rest_translation.jpg": "Translational Movement",
    "mean_mc_func.nii.gz": "Mean of Motion Corrected 4D Data",
    "initial_mean_mc_func.nii.gz": "Initial Mask on Mean Functional Image",
    "mask_mean_mc_func.nii.gz": "Final Mask on Mean Functional Image",
    "cleaned_mc_func.nii.gz": "Cleaned Functional Image (pre-Bias Correction)",
    "cleaned_N4_mean_mc_func.nii.gz": "Cleaned Mean Functional Image",
    "before_despiking_spikecountTC.png": "No of spikes before despiking data",
    "after_despiking_spikecountTC.png": "No of spikes after despiking data",
    "despike_cleaned_mc_func.nii.gz": "Despiked Functional Data",
    "cleaned_sm_despike_cleaned_mc_func.nii.gz": "1 Voxel Smoothing after Despiking"
}

section_map = {
    "Raw": {
        "files": ["G1_cp.nii.gz"],
        "color": (128, 0, 128)
    },
    "Motion Correction": {
        "files": ["mc_func.nii.gz", "rest_rotation.jpg", "rest_translation.jpg", "mean_mc_func.nii.gz"],
        "color": (0, 102, 204)
    },
    "Masking": {
        "files": ["initial_mean_mc_func.nii.gz", "mask_mean_mc_func.nii.gz"],
        "color": (0, 153, 76)
    },
    "Bias Field Correction": {
        "files": ["cleaned_mc_func.nii.gz", "cleaned_N4_mean_mc_func.nii.gz"],
        "color": (255, 153, 0)
    },
    "Despiking": {
        "files": [
            "before_despiking_spikecountTC.png", "after_despiking_spikecountTC.png",
            "despike_cleaned_mc_func.nii.gz", "cleaned_sm_despike_cleaned_mc_func.nii.gz"
        ],
        "color": (204, 0, 0)
    }
}

HEADER_COLOR = (0, 51, 102)  # Navy blue header for all

def get_section_and_color(filename):
    for section, details in section_map.items():
        if filename in details["files"]:
            return section, details["color"]
    return "Other", (100, 100, 100)

# ========== UTILITIES ==========

def save_all_slices_with_label(nifti_path, base_outpath, average_4d=False):
    img = nib.load(nifti_path)
    data = img.get_fdata()
    if data.ndim == 4 and average_4d:
        data = np.mean(data[..., :min(100, data.shape[3])], axis=3)
    if data.ndim == 4:
        data = data[..., 0]
    os.makedirs(os.path.dirname(base_outpath), exist_ok=True)
    saved_images = []
    for i in range(data.shape[2]):
        out_path = f"{base_outpath}_slice{i:02d}.png"
        plt.imshow(data[:, :, i].T, cmap="gray", origin="lower")
        plt.axis("off")
        plt.text(5, 5, f"Slice {i}", color='red', fontsize=8, ha='left', va='top',
                 bbox=dict(facecolor='white', edgecolor='none', alpha=0.7))
        plt.savefig(out_path, bbox_inches="tight", pad_inches=0)
        plt.close()
        saved_images.append((out_path, i))
    return saved_images

def add_textbox(slide, text, left, top, font_size=14, bold=False, italic=False, color=(0, 0, 0)):
    box = slide.shapes.add_textbox(left, top, Inches(8), Inches(0.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = "Calibri"
    font.size = Pt(font_size)
    font.bold = bold
    font.italic = italic
    font.color.rgb = RGBColor(*color)

def add_header_bar(slide, title):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(*HEADER_COLOR)
    bar.line.fill.background()
    add_textbox(slide, title, Inches(0.4), Inches(0.1), font_size=18, bold=True, color=(255, 255, 255))

def add_footer_bar(slide, section, section_color, timestamp, slide_num):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), prs.slide_height - Inches(0.4), prs.slide_width, Inches(0.4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(*section_color)
    bar.line.fill.background()
    add_textbox(slide, timestamp, Inches(0.3), prs.slide_height - Inches(0.35), font_size=10, italic=True, color=(255, 255, 255))
    add_textbox(slide, section, prs.slide_width / 2 - Inches(1), prs.slide_height - Inches(0.35), font_size=10, bold=True, color=(255, 255, 255))
    add_textbox(slide, f"Slide {slide_num}", prs.slide_width - Inches(1.5), prs.slide_height - Inches(0.35), font_size=10, color=(255, 255, 255))

# ========== MAIN SCRIPT ==========

prs = Presentation()
prs.slide_width = Inches(17.77)
prs.slide_height = Inches(10)
timestamp = datetime.now().strftime("Generated on %Y-%m-%d at %H:%M:%S")

for fname in tqdm(file_list, desc="Creating slides"):
    file_path = os.path.join(folder_path, fname)
    title = title_map.get(fname, fname)
    section, section_color = get_section_and_color(fname)
    is_nifti = fname.endswith(".nii.gz")
    is_4d = fname in [
        "G1_cp.nii.gz", "mc_func.nii.gz",
        "despike_cleaned_mc_func.nii.gz",
        "cleaned_sm_despike_cleaned_mc_func.nii.gz"
    ]

    if not os.path.exists(file_path):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide_num = len(prs.slides)
        add_header_bar(slide, title)
        add_footer_bar(slide, section, section_color, timestamp, slide_num)
        add_textbox(slide, f"{title} (File not found)", Inches(0.5), Inches(1.0), font_size=14, color=(255, 0, 0))
        continue

    if is_nifti:
        base_outpath = os.path.join(temp_dir, fname.replace(".nii.gz", ""))
        slice_imgs = save_all_slices_with_label(file_path, base_outpath, average_4d=is_4d)
        for i in range(0, len(slice_imgs), 16):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide_num = len(prs.slides)
            add_header_bar(slide, title + f" (Slices {i+1}-{min(i+16, len(slice_imgs))})")
            add_footer_bar(slide, section, section_color, timestamp, slide_num)
            for idx, (img_path, slice_num) in enumerate(slice_imgs[i:i+16]):
                row, col = divmod(idx, 4)
                x = Inches(0.3 + col * 1.75)
                y = Inches(1.0 + row * 1.75)
                slide.shapes.add_picture(img_path, x, y, height=Inches(1.6))
    else:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide_num = len(prs.slides)
        add_header_bar(slide, title)
        add_footer_bar(slide, section, section_color, timestamp, slide_num)
        slide.shapes.add_picture(file_path, Inches(1), Inches(1.5), height=Inches(5.5))

prs.save(output_pptx)
print(f"\n✅ PowerPoint saved to: {output_pptx}")
