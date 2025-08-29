import os
import subprocess
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from datetime import datetime
from tqdm import tqdm
import glob
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
import sys

# ========== User-defined paths ==========

output_dir = sys.argv[1]
anatomy_path = sys.argv[2]
overlay_path = sys.argv[3]
split_dir = sys.argv[4]
# 
# output_dir = "/Users/njain/Desktop/renders"
# anatomy_path = "/Users/njain/Desktop/RGRO_250430_0224_RN_SD_020/11anatomy/anatomy.nii.gz"
# overlay_path = "/Users/njain/Desktop/RGRO_250430_0224_RN_SD_020/10functionalEPI/Coregistered_SCM.nii.gz"
# split_dir = "/Users/njain/Desktop/RGRO_250430_0224_RN_SD_020/10functionalEPI/split_volumes"  # Make sure this is populated with volXXXX.nii.gz
# ========================================

# PowerPoint file with timestamp
pptx_path = f"/Users/njain/Desktop/render_output_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
os.makedirs(output_dir, exist_ok=True)
prs = Presentation()

slice_indices = [22,23,24,25]  # For 2×2 grid
volumes = sorted(glob.glob(os.path.join(split_dir, "vol*.nii.gz")))
overlay_filename = os.path.basename(overlay_path)
timestamp_str = datetime.now().strftime("Generated on: %Y-%m-%d %H:%M")

# Grid layout (2×2)
positions = [
    (Inches(0.5), Inches(1.0)),  # top-left
    (Inches(4.8), Inches(1.0)),  # top-right
    (Inches(0.5), Inches(3.8)),  # bottom-left
    (Inches(4.8), Inches(3.8)),  # bottom-right
]

# Main loop
for vol_index, vol_path in enumerate(tqdm(volumes, desc="Rendering volumes")):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title at top
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = f"Volume: {vol_index}"
    title_frame.paragraphs[0].font.size = Pt(20)
    title_frame.paragraphs[0].font.bold = True

    # Loop over slices
    for idx, z in enumerate(slice_indices):
        out_img = os.path.join(output_dir, f"vol{vol_index:04d}_z{z}.png")

        cmd = [
            "render",
            "-of", out_img,
            "-sz", "1600", "1200",
            "-c", "10",
            anatomy_path,
                "-n", "anatomy",
                "-dr", "0", "13419.91",
                "-cr", "0", "13419.91",
                "-cm", "greyscale",
                "-b", "49.75", "-c", "49.90", "-a", "100",
            vol_path,
                "-n", f"Volume_{vol_index}",
                "-dr", "0", "40.0",
                "-cr", "0", "40.0",
                "-cm", "brain_colours_blackbdy",
                "-b", "33.71", "-c", "72.30",
                "-ma", "-mr", "0", "25", "-a", "100"
        ]

        subprocess.run(cmd, check=True)

        # Place image
        left, top = positions[idx]
        img = slide.shapes.add_picture(out_img, left, top, width=Inches(4), height=Inches(2.6))

        # Add z-slice label inside image, bottom-left corner
        label_box = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(2.3), Inches(2), Inches(0.3))
        label_frame = label_box.text_frame
        label_frame.text = f"Z-slice: {z}"
        label_frame.paragraphs[0].font.size = Pt(11)
        label_frame.paragraphs[0].font.bold = True
        label_frame.paragraphs[0].font.color.rgb = RGBColor(230, 230, 230)  # Light inside-label
        label_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT 

    # Add footer (bottom-left)
    footer_text = f"{timestamp_str} | Overlay file: {overlay_filename}"
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(7.5), Inches(0.3))
    footer_frame = footer_box.text_frame
    footer_frame.text = footer_text
    footer_frame.paragraphs[0].font.size = Pt(10)
    footer_frame.paragraphs[0].font.italic = True
    footer_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)  # Gray

# Save
prs.save(pptx_path)
print(f"✅ Saved PowerPoint with 2×2 grid, slice labels, and footer: {pptx_path}")


# import os
# import subprocess
# from pptx import Presentation
# from pptx.util import Inches, Pt
# from pptx.dml.color import RGBColor
# from datetime import datetime
# from tqdm import tqdm
# import glob

# # Configuration
# output_dir = "/Users/njain/Desktop/renders"
# pptx_path = f"/Users/njain/Desktop/render_output_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
# anatomy_path = "/Users/njain/Desktop/RGRO_250430_0224_RN_SD_020/11anatomy/anatomy.nii.gz"
# split_dir = "/Users/njain/Desktop/RGRO_250430_0224_RN_SD_020/10functionalEPI/split_volumes"  # Make sure this is populated with volXXXX.nii.gz

# # Ensure output dir exists
# os.makedirs(output_dir, exist_ok=True)

# # PowerPoint setup
# prs = Presentation()
# slice_indices = [21, 23, 24]  # 4 slices for 2x2 grid

# # Get sorted volume paths
# volumes = sorted(glob.glob(os.path.join(split_dir, "vol*.nii.gz")))

# # Grid layout positions
# positions = [
#     (Inches(0.5), Inches(1.0)),  # top-left
#     (Inches(4.8), Inches(1.0)),  # top-right
#     (Inches(0.5), Inches(3.8)),  # bottom-left
#     (Inches(4.8), Inches(3.8)),  # bottom-right
# ]

# # Loop through each volume
# for vol_index, vol_path in enumerate(tqdm(volumes, desc="Rendering volumes")):
#     slide = prs.slides.add_slide(prs.slide_layouts[6])

#     # Title: Volume X
#     title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8), Inches(0.6))
#     title_frame = title_box.text_frame
#     title_frame.text = f"Volume: {vol_index}"
#     title_frame.paragraphs[0].font.size = Pt(20)
#     title_frame.paragraphs[0].font.bold = True

#     # Render and place each slice in 2x2 grid
#     for idx, z in enumerate(slice_indices):
#         out_img = os.path.join(output_dir, f"vol{vol_index:04d}_z{z}.png")

#         cmd = [
#             "render",
#             "-of", out_img,
#             "-sz", "1600", "1200",
#             "-c", "10",
#             anatomy_path,
#                 "-n", "anatomy",
#                 "-dr", "0", "13419.91",
#                 "-cr", "0", "13419.91",
#                 "-cm", "greyscale",
#                 "-b", "49.75", "-c", "49.90", "-a", "100",
#             vol_path,
#                 "-n", f"Volume_{vol_index}",
#                 "-dr", "0", "8.0",
#                 "-cr", "0", "8.0",
#                 "-cm", "brain_colours_blackbdy",
#                 "-b", "33.71", "-c", "72.30",
#                 "-ma", "-mr", "0", "6", "-a", "100"
#         ]

#         subprocess.run(cmd, check=True)

#         # Add image to slide in grid layout
#         left, top = positions[idx]
#         slide.shapes.add_picture(out_img, left, top, width=Inches(4), height=Inches(2.6))

#     # Add timestamp footer
#     timestamp = datetime.now().strftime("Generated on: %Y-%m-%d %H:%M")
#     footer_box = slide.shapes.add_textbox(Inches(6.5), Inches(6.8), Inches(3), Inches(0.3))
#     footer_frame = footer_box.text_frame
#     footer_frame.text = timestamp
#     footer_frame.paragraphs[0].font.size = Pt(10)
#     footer_frame.paragraphs[0].font.italic = True
#     footer_frame.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)  # Light gray

# # Save presentation
# prs.save(pptx_path)
# print(f"✅ Saved PowerPoint with 2×2 grid per volume: {pptx_path}")



