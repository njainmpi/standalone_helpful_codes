import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image
from datetime import datetime
from math import ceil

# --- CONFIG ---
TITLE_FONT_SIZE = Pt(28)
TITLE_COLOR = RGBColor(0x00, 0x33, 0x66)  # Dark blue
TITLE_BOLD = True
SLIDE_WIDTH = 10.0
SLIDE_HEIGHT = 7.5
CHARS_PER_LINE = 60

# --- UTILS ---
def title_from_filename(filename):
    name = os.path.splitext(os.path.basename(filename))[0]
    return name.replace("_", " ").replace("-", " ").title()

def estimate_title_height(title_text, font_size_pt=28, chars_per_line=60):
    line_count = ceil(len(title_text) / chars_per_line)
    line_height_inch = font_size_pt * 1.2 / 72  # 1.2 line spacing
    return line_count * line_height_inch

def add_footer(slide, timestamp):
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(9), Inches(0.3))
    text_frame = textbox.text_frame
    text_frame.text = f"Generated on {timestamp}"
    text_frame.paragraphs[0].font.size = Pt(10)

def add_title(slide, title_text):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    font = run.font
    font.size = TITLE_FONT_SIZE
    font.bold = TITLE_BOLD
    font.color.rgb = TITLE_COLOR

def add_image_centered(slide, img_path, title_text, font_size_pt=28):
    top_offset = estimate_title_height(title_text, font_size_pt)
    max_width = SLIDE_WIDTH
    max_height = SLIDE_HEIGHT - top_offset

    with Image.open(img_path) as img:
        width_px, height_px = img.size
        img_ratio = width_px / height_px
        box_ratio = max_width / max_height

        if img_ratio > box_ratio:
            pic_width = Inches(max_width)
            pic_height = Inches(max_width / img_ratio)
        else:
            pic_height = Inches(max_height)
            pic_width = Inches(max_height * img_ratio)

        left = Inches((SLIDE_WIDTH - pic_width.inches) / 2)
        top = Inches(top_offset + ((SLIDE_HEIGHT - top_offset - pic_height.inches) / 2))

        slide.shapes.add_picture(img_path, left, top, width=pic_width, height=pic_height)

# --- MAIN ---
if len(sys.argv) != 3:
    print("Usage: python3 auto_ppt_from_folder.py /path/to/image_folder /path/to/output_name.pptx")
    sys.exit(1)

image_dir = sys.argv[1]
pptx_path = sys.argv[2]

# Create new or open existing presentation
prs = Presentation() if not os.path.exists(pptx_path) else Presentation(pptx_path)

# Use a blank layout (no placeholders)
blank_layout = next((layout for layout in prs.slide_layouts if len(layout.placeholders) == 0), prs.slide_layouts[6])

# Supported image formats
extensions = (".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp", ".gif", ".webp")
image_files = sorted([
    os.path.join(image_dir, f)
    for f in os.listdir(image_dir)
    if f.lower().endswith(extensions)
])

if not image_files:
    print("❌ No supported image files found in:", image_dir)
    sys.exit(1)

timestamp = datetime.now().strftime("%Y-%m-%d %H:%M")

for img_path in image_files:
    title = title_from_filename(img_path)

    try:
        with Image.open(img_path) as img:
            img.verify()
    except Exception as e:
        print(f"⚠️ Skipping {img_path}: {e}")
        continue

    blank_layout = next((layout for layout in prs.slide_layouts if len(layout.placeholders) == 0), prs.slide_layouts[6])
    slide = prs.slides.add_slide(blank_layout)

    add_title(slide, title)
    add_image_centered(slide, img_path, title)
    add_footer(slide, timestamp)

    print(f"✅ Added slide: {title}")

prs.save(pptx_path)
print(f"\n🎉 Presentation saved to: {pptx_path}")



    
